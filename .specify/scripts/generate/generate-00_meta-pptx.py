#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define color scheme (Banking/Professional theme)
COLOR_PRIMARY = RGBColor(46, 80, 144)  # Dark Blue #2E5090
COLOR_SECONDARY = RGBColor(68, 114, 196)  # Medium Blue #4472C4
COLOR_ACCENT = RGBColor(91, 155, 213)  # Light Blue #5B9BD5
COLOR_TEXT = RGBColor(51, 51, 51)  # Dark Gray #333333
COLOR_BG = RGBColor(245, 245, 245)  # Light Gray Background

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5),
        Inches(8), Inches(1.5)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = COLOR_PRIMARY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets):
    """Add a content slide with bullets"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.9)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_PRIMARY
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Content area
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3),
        Inches(8.6), Inches(5.7)
    )
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i > 0:
            p = text_frame.add_paragraph()
        else:
            p = text_frame.paragraphs[0]
        p.text = bullet
        p.level = 0
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6)
        p.space_after = Pt(6)

    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    background.line.fill.background()

    # Title bar
    title_bar = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.9)
    )
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = COLOR_PRIMARY
    title_bar.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(9), Inches(0.6)
    )
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(32)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Table
    table = slide.shapes.add_table(
        len(rows) + 1, len(headers),
        Inches(0.5), Inches(1.5),
        Inches(9), Inches(5.5)
    ).table

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(213, 232, 240)  # Light blue
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    # Data rows
    for i, row in enumerate(rows):
        for j, value in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(value)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = COLOR_TEXT

    return slide

# Slide 1: Title
add_title_slide(
    prs,
    "客戶洗錢風險事件摘要模型",
    "專案基本資料 (00_meta) | 2025-11-14"
)

# Slide 2: 專案總覽
add_content_slide(
    prs,
    "專案總覽",
    [
        "專案代號: RISK-AML-MDL-001",
        "專案名稱: 客戶洗錢風險事件摘要模型",
        "發起單位: 法遵部 / 風控部",
        "專案類型: 新系統開發",
        "專案性質: 監理申報 / 法遵要求 / 風險控制",
        "優先等級: P0-Critical"
    ]
)

# Slide 3: 關鍵時程
add_table_slide(
    prs,
    "關鍵時程",
    ["里程碑", "預定日期", "狀態"],
    [
        ["需求確認完成", "2025-12-15", "規劃中"],
        ["設計審查通過", "2026-01-31", "規劃中"],
        ["開發完成", "2026-04-15", "規劃中"],
        ["UAT 測試完成", "2026-05-15", "規劃中"],
        ["正式上線", "2026-05-31", "規劃中"]
    ]
)

# Slide 4: 預期效益
add_content_slide(
    prs,
    "預期效益",
    [
        "量化效益:",
        "  • 減少人工風險分析時間: 預估節省 60% 以上",
        "  • 提升洗錢風險識別準確率: 目標提升至 85% 以上",
        "  • 加速風險事件摘要產出: 從 3 天縮短至即時生成",
        "",
        "質化效益:",
        "  • 符合金融監理機關反洗錢法規要求",
        "  • 強化客戶風險管控能力，降低法遵風險",
        "  • 提供管理層即時風險儀表板，支援決策"
    ]
)

# Slide 5: 專案背景 - 問題陳述
add_content_slide(
    prs,
    "問題陳述 - 現行痛點",
    [
        "客戶洗錢風險事件分散於多個系統，缺乏整合性摘要",
        "人工彙整風險事件耗時且容易遺漏關鍵資訊",
        "風險評分標準不一致，依賴人工判斷，缺乏客觀性",
        "無法即時掌握高風險客戶動態，影響預警能力",
        "監理申報準備作業繁複，資料品質難以保證"
    ]
)

# Slide 6: 業務驅動因素
add_content_slide(
    prs,
    "業務驅動因素",
    [
        "法規要求:",
        "  • 遵循《洗錢防制法》及相關子法規定",
        "  • 滿足金融監理機關（如金管會）監理要求",
        "  • 配合國際反洗錢標準（FATF 建議）",
        "",
        "風險管理需求:",
        "  • 強化客戶風險分級管理機制",
        "  • 建立可量化、可追溯的風險評估模型",
        "  • 提升異常交易偵測與預警能力"
    ]
)

# Slide 7: 專案範圍 - 範圍內
add_content_slide(
    prs,
    "專案範圍 - In Scope",
    [
        "客戶風險事件資料整合引擎",
        "洗錢風險評分模型（基於規則與機器學習）",
        "客戶風險事件摘要自動產生功能",
        "風險等級分類與標註（高/中/低風險）",
        "風險儀表板與視覺化介面",
        "風險趨勢分析與統計報表",
        "高風險客戶自動預警通知",
        "監理申報資料匯出功能"
    ]
)

# Slide 8: 專案範圍 - 範圍外
add_content_slide(
    prs,
    "專案範圍 - Out of Scope",
    [
        "不包含實時交易監控功能（由既有系統負責）",
        "不包含客戶盡職調查(KYC)流程管理",
        "不包含可疑交易申報(STR/SAR)的案件調查流程",
        "不包含帳戶凍結或交易阻擋功能",
        "不涉及客戶身分驗證(eKYC)機制",
        "不包含反洗錢教育訓練管理",
        "第一期不包含跨境交易風險分析"
    ]
)

# Slide 9: 高階風險
add_table_slide(
    prs,
    "高階風險",
    ["風險ID", "風險描述", "機率", "影響"],
    [
        ["R001", "資料品質不佳導致模型準確度低", "高", "高"],
        ["R002", "跨系統資料整合複雜度高於預期", "中", "高"],
        ["R003", "模型可解釋性不足無法通過稽核", "中", "高"],
        ["R004", "關鍵人力（資料科學家）取得困難", "中", "高"]
    ]
)

# Slide 10: 成功標準
add_content_slide(
    prs,
    "成功標準",
    [
        "完成標準 (Definition of Done):",
        "  • 所有範圍內功能開發完成並通過測試",
        "  • 風險模型準確率達到設定門檻（目標 ≥ 85%）",
        "  • UAT 測試通過率 ≥ 95%",
        "  • 資安測試通過（無高風險與中風險弱點）",
        "",
        "驗收標準 (Acceptance Criteria):",
        "  • 法遵部與風控部簽核業務驗收",
        "  • 模型驗證委員會通過模型驗證",
        "  • 資安檢測通過（弱點掃描、滲透測試）"
    ]
)

# Save presentation
prs.save("/Users/ro9air/SDD_repo/output/00_meta.pptx")
print("PPTX presentation generated successfully: /Users/ro9air/SDD_repo/output/00_meta.pptx")
